from dataclasses import dataclass
from io import BytesIO

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas


@dataclass(frozen=True)
class ParserFixture:
    media_type: str
    content: bytes
    expected_kinds: tuple[str, ...]


def searchable_pdf_fixture(*page_texts: str) -> bytes:
    output = BytesIO()
    document = canvas.Canvas(output)
    if "STSong-Light" not in pdfmetrics.getRegisteredFontNames():
        pdfmetrics.registerFont(  # pyright: ignore[reportUnknownMemberType]
            UnicodeCIDFont("STSong-Light")
        )
    for text in page_texts:
        font = "STSong-Light" if any(ord(character) > 127 for character in text) else "Helvetica"
        document.setFont(font, 12)
        if text:
            document.drawString(72, 720, text)
        document.showPage()
    document.save()
    return output.getvalue()


def docx_fixture() -> bytes:
    document = Document()
    document.add_heading("Overview", level=1)
    paragraph = document.add_paragraph()
    paragraph.add_run("Hello ")
    paragraph.add_run("世界")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "name"
    table.cell(0, 1).text = "value"
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def empty_docx_fixture() -> bytes:
    document = Document()
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def pptx_fixture() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    text_box.text_frame.text = "Body 世界"
    notes_text_frame = slide.notes_slide.notes_text_frame
    assert notes_text_frame is not None
    notes_text_frame.text = "Speaker notes"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def empty_pptx_fixture() -> bytes:
    presentation = Presentation()
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def xlsx_fixture() -> bytes:
    workbook = Workbook()
    worksheet = workbook.active
    assert worksheet is not None
    worksheet.title = "数据"
    worksheet.append(["name", "value"])
    worksheet.append(["中文", 1])
    output = BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def empty_xlsx_fixture() -> bytes:
    workbook = Workbook()
    output = BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def parser_contract_fixtures() -> tuple[ParserFixture, ...]:
    return (
        ParserFixture(
            media_type="text/plain",
            content="Hello\r\n世界\x00\x01\x7f\x80\x9f".encode(),
            expected_kinds=("text",),
        ),
        ParserFixture(
            media_type="text/markdown",
            content="# 标题\r\n\r\nBody".encode(),
            expected_kinds=("heading", "paragraph"),
        ),
        ParserFixture(
            media_type="text/csv",
            content="name,value\r\n中文,1".encode(),
            expected_kinds=("sheet_rows",),
        ),
        ParserFixture(
            media_type="text/html",
            content=b"<h1>Title</h1><p>Hello</p>",
            expected_kinds=("heading", "paragraph"),
        ),
        ParserFixture(
            media_type="application/pdf",
            content=searchable_pdf_fixture("English", "中文"),
            expected_kinds=("text", "text"),
        ),
        ParserFixture(
            media_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            content=docx_fixture(),
            expected_kinds=("heading", "paragraph", "table"),
        ),
        ParserFixture(
            media_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            content=pptx_fixture(),
            expected_kinds=("slide", "slide"),
        ),
        ParserFixture(
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            content=xlsx_fixture(),
            expected_kinds=("sheet_rows",),
        ),
    )


def whitespace_parser_fixtures() -> tuple[tuple[str, bytes], ...]:
    return (
        ("text/plain", b" \r\n\t"),
        ("text/markdown", b" \r\n\t"),
        ("text/csv", b" \r\n\t"),
        ("text/html", b"<html><body> \r\n\t</body></html>"),
        ("application/pdf", searchable_pdf_fixture("")),
        (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            empty_docx_fixture(),
        ),
        (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            empty_pptx_fixture(),
        ),
        (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            empty_xlsx_fixture(),
        ),
    )
