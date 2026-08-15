from collections.abc import Callable
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from cairn_worker.errors import WorkerFailure
from cairn_worker.parsers import BlockKind
from cairn_worker.parsers.pptx import PptxParser
from pptx import Presentation
from pptx.util import Inches

from apps.worker.tests.fixture_factory import empty_pptx_fixture


def _structured_pptx() -> bytes:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    first_text = first.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    first_text.text_frame.text = "First 世界"
    table = first.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(2)).table
    table.cell(0, 0).text = "name"
    table.cell(0, 1).text = "value"
    table.cell(1, 0).text = "alpha"
    table.cell(1, 1).text = "一"
    group = first.shapes.add_group_shape()
    grouped_text = group.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
    grouped_text.text_frame.text = "Grouped"
    notes_text_frame = first.notes_slide.notes_text_frame
    assert notes_text_frame is not None
    notes_text_frame.text = "Speaker notes"

    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second_text = second.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    second_text.text_frame.text = "Second"

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _presentation_with_external_hyperlink() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = text_box.text_frame.paragraphs[0].add_run()
    run.text = "Visible link label"
    run.hyperlink.address = "https://127.0.0.1:1/private"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _append_member(package: bytes, name: str, content: bytes) -> bytes:
    output = BytesIO()
    with ZipFile(BytesIO(package), "r") as source, ZipFile(
        output, "w", compression=ZIP_DEFLATED
    ) as target:
        for member in source.infolist():
            target.writestr(member, source.read(member))
        target.writestr(name, content)
    return output.getvalue()


def _replace_member(package: bytes, name: str, transform: Callable[[bytes], bytes]) -> bytes:
    output = BytesIO()
    with ZipFile(BytesIO(package), "r") as source, ZipFile(
        output, "w", compression=ZIP_DEFLATED
    ) as target:
        for member in source.infolist():
            data = source.read(member)
            if member.filename == name:
                data = transform(data)
            target.writestr(member, data)
    return output.getvalue()


def test_pptx_parser_emits_body_and_notes_once_in_slide_order() -> None:
    """Break caught: body/table/group text and notes must stay ordered without metadata copies."""
    blocks = PptxParser().parse(BytesIO(_structured_pptx()))

    assert [block.kind for block in blocks] == [
        BlockKind.SLIDE,
        BlockKind.SLIDE,
        BlockKind.SLIDE,
    ]
    assert [block.text for block in blocks] == [
        "First 世界\nname\tvalue\nalpha\t一\nGrouped",
        "Speaker notes",
        "Second",
    ]
    assert [block.locator.model_dump(by_alias=True) for block in blocks] == [
        {"type": "pptx", "slide": 1, "area": "body"},
        {"type": "pptx", "slide": 1, "area": "notes"},
        {"type": "pptx", "slide": 2, "area": "body"},
    ]


def test_pptx_parser_reads_only_the_visible_label_of_an_external_hyperlink() -> None:
    """Break caught: parsing must not fetch or index an external relationship target."""
    blocks = PptxParser().parse(BytesIO(_presentation_with_external_hyperlink()))

    assert [block.text for block in blocks] == ["Visible link label"]
    assert "127.0.0.1" not in blocks[0].text


def test_pptx_parser_rejects_empty_decks_and_slides() -> None:
    """Break caught: a deck with no visible body or notes must not become ready."""
    presentation = Presentation(BytesIO(empty_pptx_fixture()))
    presentation.slides.add_slide(presentation.slide_layouts[6])
    output = BytesIO()
    presentation.save(output)

    for content in (empty_pptx_fixture(), output.getvalue()):
        with pytest.raises(WorkerFailure) as caught:
            PptxParser().parse(BytesIO(content))
        assert caught.value.code == "no_extractable_text"


def test_pptx_parser_rejects_malformed_and_macro_bearing_packages_safely() -> None:
    """Break caught: malformed XML and presentation macros must stop before content walking."""
    contents = (
        b"not a presentation",
        _append_member(empty_pptx_fixture(), "ppt/vbaProject.bin", b"private macro"),
    )
    for content in contents:
        with pytest.raises(WorkerFailure) as caught:
            PptxParser().parse(BytesIO(content))
        assert caught.value.code == "parser_failed"
        assert caught.value.retryable is False
        assert caught.value.safe_detail == "worker handler or parser failed"


def test_pptx_parser_enforces_its_slide_work_limit(monkeypatch: pytest.MonkeyPatch) -> None:
    """Break caught: slide traversal must stop at a deterministic parser work boundary."""
    from cairn_worker.parsers import pptx as pptx_parser

    monkeypatch.setattr(pptx_parser, "PPTX_MAX_SLIDES", 1)
    with pytest.raises(WorkerFailure) as caught:
        PptxParser().parse(BytesIO(_structured_pptx()))
    assert caught.value.code == "parser_failed"


def test_pptx_parser_excludes_hidden_slides_shapes_groups_and_notes() -> None:
    package = _structured_pptx()
    hidden_slide = _replace_member(
        package,
        "ppt/slides/slide1.xml",
        lambda xml: xml.replace(b"<p:sld ", b'<p:sld show="0" ', 1),
    )
    blocks = PptxParser().parse(BytesIO(hidden_slide))
    assert [block.text for block in blocks] == ["Second"]
    assert blocks[0].locator.model_dump(by_alias=True)["slide"] == 2

    def hide_selected(xml: bytes) -> bytes:
        xml = xml.replace(b'<p:cNvPr id="2"', b'<p:cNvPr hidden="1" id="2"', 1)
        xml = xml.replace(b'<p:cNvPr id="4"', b'<p:cNvPr hidden="true" id="4"', 1)
        return xml

    hidden_shapes = _replace_member(package, "ppt/slides/slide1.xml", hide_selected)
    hidden_shapes = _replace_member(
        hidden_shapes,
        "ppt/notesSlides/notesSlide1.xml",
        lambda xml: xml.replace(
            b'<p:cNvPr id="3"', b'<p:cNvPr hidden="1" id="3"', 1
        ),
    )
    blocks = PptxParser().parse(BytesIO(hidden_shapes))
    texts = [block.text for block in blocks]
    assert "First 世界" not in "\n".join(texts)
    assert "Grouped" not in "\n".join(texts)
    assert "Speaker notes" not in "\n".join(texts)
    assert "name\tvalue\nalpha\t一" in texts[0]
    assert texts[-1] == "Second"

    hidden_table = _replace_member(
        package,
        "ppt/slides/slide1.xml",
        lambda xml: xml.replace(b'<p:cNvPr id="3"', b'<p:cNvPr hidden="on" id="3"', 1),
    )
    table_texts = "\n".join(block.text for block in PptxParser().parse(BytesIO(hidden_table)))
    assert "name" not in table_texts
    assert "First 世界" in table_texts


def test_pptx_parser_rejects_xml_work_before_presentation_construction(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from cairn_worker.parsers import office_safety
    from cairn_worker.parsers import pptx as pptx_parser

    constructed = [False]
    monkeypatch.setattr(office_safety, "OPC_MAX_XML_TEXT_CHARACTERS", 1)

    def unexpected_presentation(*args: object, **kwargs: object) -> object:
        del args, kwargs
        constructed[0] = True
        raise AssertionError

    monkeypatch.setattr(pptx_parser, "Presentation", unexpected_presentation)
    with pytest.raises(WorkerFailure) as caught:
        PptxParser().parse(BytesIO(_structured_pptx()))
    assert caught.value.code == "parser_failed"
    assert constructed[0] is False


def test_pptx_parser_runtime_shape_work_has_exact_boundary(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from cairn_worker.parsers import pptx as pptx_parser

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1)).text = "one"
    output = BytesIO()
    presentation.save(output)
    monkeypatch.setattr(pptx_parser, "PPTX_MAX_SHAPE_AND_CELL_WORK", 1)
    PptxParser().parse(BytesIO(output.getvalue()))

    slide.shapes.add_textbox(Inches(2), Inches(1), Inches(1), Inches(1)).text = "two"
    output = BytesIO()
    presentation.save(output)
    with pytest.raises(WorkerFailure) as caught:
        PptxParser().parse(BytesIO(output.getvalue()))
    assert caught.value.code == "parser_failed"


def test_pptx_parser_runtime_table_cell_work_has_exact_boundary(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from cairn_worker.parsers import pptx as pptx_parser

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(2), Inches(1)).table
    table.cell(0, 0).text = "one"
    output = BytesIO()
    presentation.save(output)
    monkeypatch.setattr(pptx_parser, "PPTX_MAX_SHAPE_AND_CELL_WORK", 2)
    PptxParser().parse(BytesIO(output.getvalue()))

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 1, Inches(1), Inches(1), Inches(2), Inches(1)).table
    table.cell(0, 0).text = "one"
    table.cell(1, 0).text = "two"
    output = BytesIO()
    presentation.save(output)
    with pytest.raises(WorkerFailure) as caught:
        PptxParser().parse(BytesIO(output.getvalue()))
    assert caught.value.code == "parser_failed"


def test_pptx_parser_normalizes_whitespace_around_hidden_and_show_flags() -> None:
    package = _structured_pptx()
    hidden_slide = _replace_member(
        package,
        "ppt/slides/slide1.xml",
        lambda xml: xml.replace(b"<p:sld ", b'<p:sld show=" 0 \t" ', 1),
    )
    assert [block.text for block in PptxParser().parse(BytesIO(hidden_slide))] == ["Second"]

    def hide_body(xml: bytes) -> bytes:
        xml = xml.replace(b'<p:cNvPr id="2"', b'<p:cNvPr hidden=" true \n" id="2"', 1)
        xml = xml.replace(b'<p:cNvPr id="3"', b'<p:cNvPr hidden=" on " id="3"', 1)
        xml = xml.replace(b'<p:cNvPr id="4"', b'<p:cNvPr hidden=" yes " id="4"', 1)
        return xml

    hidden = _replace_member(package, "ppt/slides/slide1.xml", hide_body)
    hidden = _replace_member(
        hidden,
        "ppt/notesSlides/notesSlide1.xml",
        lambda xml: xml.replace(
            b'<p:cNvPr id="3"', b'<p:cNvPr hidden=" 1 " id="3"', 1
        ),
    )
    assert [block.text for block in PptxParser().parse(BytesIO(hidden))] == ["Second"]
