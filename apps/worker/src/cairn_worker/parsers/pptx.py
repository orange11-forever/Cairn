from io import BytesIO
from typing import Any, BinaryIO, Protocol, cast

from cairn_api.knowledge.schemas import PptxLocator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.table import Table
from pptx.text.text import TextFrame

from cairn_worker.parsers import BlockKind, DocumentParser, ParsedBlock, read_parser_source
from cairn_worker.parsers.limits import MAX_PARSED_BLOCKS, ensure_block_capacity
from cairn_worker.parsers.office_safety import validate_opc_package

PPTX_MAX_SLIDES = MAX_PARSED_BLOCKS // 2
PPTX_MAX_SHAPE_AND_CELL_WORK = 1_000_000
_TRUTHY_XML = {"1", "true", "on", "yes"}


class _TextShape(Protocol):
    @property
    def text_frame(self) -> TextFrame: ...


class _TableShape(Protocol):
    @property
    def table(self) -> Table: ...


def _xml_hidden(element: object) -> bool:
    candidates = cast(Any, element).xpath(".//p:cNvPr")
    if not candidates:
        return False
    value = candidates[0].get("hidden")
    return value is not None and value.strip().casefold() in _TRUTHY_XML


def _text_frame_text(text_frame: TextFrame) -> str:
    return "\n".join(
        text
        for paragraph in text_frame.paragraphs
        if (text := paragraph.text.strip())
    )


def _table_text(table: Table, work: list[int]) -> str:
    rows: list[str] = []
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            work[0] += 1
            if work[0] > PPTX_MAX_SHAPE_AND_CELL_WORK:
                raise ValueError("PPTX shapes exceed parser work limit")
            cells.append(_text_frame_text(cell.text_frame))
        if any(cell.strip() for cell in cells):
            rows.append("\t".join(cells))
    return "\n".join(rows)


def _shape_segments(shape: BaseShape, work: list[int]) -> list[str]:
    work[0] += 1
    if work[0] > PPTX_MAX_SHAPE_AND_CELL_WORK:
        raise ValueError("PPTX shapes exceed parser work limit")
    if _xml_hidden(shape.element):
        return []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group = cast(GroupShape, shape)
        return [
            segment
            for child in group.shapes
            for segment in _shape_segments(child, work)
        ]
    if shape.has_table:
        text = _table_text(cast(_TableShape, shape).table, work)
        return [text] if text.strip() else []
    if shape.has_text_frame:
        text = _text_frame_text(cast(_TextShape, shape).text_frame)
        return [text] if text.strip() else []
    return []


class PptxParser(DocumentParser):
    def _parse(self, source: BinaryIO) -> list[ParsedBlock]:
        content = read_parser_source(source)
        validate_opc_package(content, required_member="ppt/presentation.xml")
        presentation = Presentation(BytesIO(content))
        if len(presentation.slides) > PPTX_MAX_SLIDES:
            raise ValueError("PPTX slide count exceeds parser work limit")

        blocks: list[ParsedBlock] = []
        work = [0]
        for slide_number, slide in enumerate(presentation.slides, start=1):
            show = cast(str | None, cast(Any, slide.element).get("show"))
            if show is not None and show.strip().casefold() in {"0", "false", "off", "no"}:
                continue
            body = "\n".join(
                segment
                for shape in slide.shapes
                for segment in _shape_segments(shape, work)
            )
            if body.strip():
                ensure_block_capacity(len(blocks))
                blocks.append(
                    ParsedBlock(
                        kind=BlockKind.SLIDE,
                        text=body,
                        locator=PptxLocator(slide=slide_number, area="body"),
                    )
                )

            if slide.has_notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                notes = (
                    _text_frame_text(notes_text_frame)
                    if notes_text_frame is not None
                    and not _xml_hidden(
                        cast(Any, notes_text_frame._txBody).getparent()  # pyright: ignore[reportPrivateUsage]
                    )
                    else ""
                )
                if notes.strip():
                    ensure_block_capacity(len(blocks))
                    blocks.append(
                        ParsedBlock(
                            kind=BlockKind.SLIDE,
                            text=notes,
                            locator=PptxLocator(slide=slide_number, area="notes"),
                        )
                    )
        return blocks


__all__ = ["PptxParser"]
